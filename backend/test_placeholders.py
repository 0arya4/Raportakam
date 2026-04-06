from pptx import Presentation

file_path = "d:\\raportakam\\backend\\templates\\sample1.pptx"

try:
    prs = Presentation(file_path)
    print("Master Layouts available:")
    for i, layout in enumerate(prs.slide_layouts):
        text_shapes = [sh for sh in layout.shapes if hasattr(sh, "text") and sh.text]
        placeholders = list(layout.placeholders)
        print(f"Layout {i}: {layout.name} - Shapes with built-in text: {len(text_shapes)}, Placeholders: {len(placeholders)}")
        for i, ph in enumerate(placeholders):
            print(f"  - Placeholder {i} idx {ph.placeholder_format.idx} type {ph.placeholder_format.type}")
except Exception as e:
    print(f"Error: {e}")
