from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
import asyncio
import threading
import json
import os
import re
import tempfile
import subprocess
import uuid
import shutil
import time
import boto3
from botocore.config import Config
from dotenv import load_dotenv, dotenv_values
import anthropic
from system_prompt import SYSTEM_PROMPT
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

_ENV_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
load_dotenv(dotenv_path=_ENV_PATH, override=True)
_env_vals = dotenv_values(_ENV_PATH)
_anthropic_key = _env_vals.get("ANTHROPIC_API_KEY") or os.getenv("ANTHROPIC_API_KEY") or ""
print(f"[STARTUP] ENV_PATH={_ENV_PATH} key_len={len(_anthropic_key)} key_start={_anthropic_key[:15]!r}")

import logging
_logger = logging.getLogger("uvicorn.error")

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "https://raportakam.com", "https://www.raportakam.com"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("startup")
async def _startup_log():
    _logger.info(f"[STARTUP] ENV_PATH={_ENV_PATH}")
    _logger.info(f"[STARTUP] key_len={len(_anthropic_key)} key_start={_anthropic_key[:15]!r}")
    _logger.info(f"[STARTUP] key loaded OK")

claude = anthropic.Anthropic(api_key=_anthropic_key)

# ── Smart Estimation System ──────────────────────────────────────────────────
STATS_FILE = os.path.join(os.path.dirname(__file__), "generation_stats.json")
MAX_STATS = 200  # keep last 200 entries

def _load_stats():
    try:
        with open(STATS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []

def _save_stats(stats):
    try:
        with open(STATS_FILE, "w", encoding="utf-8") as f:
            json.dump(stats[-MAX_STATS:], f)
    except Exception:
        pass

def _detail_score(detail: str) -> int:
    return {"Summary": 0, "Balanced": 1, "Detailed": 2}.get(detail, 1)

def _addon_count(req) -> int:
    return sum([
        req.addon_table, req.addon_timeline, req.addon_chart_extra,
        req.addon_quotes, req.addon_comparison, req.addon_cover_page,
        req.conclusion, req.addon_references,
    ])

PRICING = {
    "claude-haiku-4-5":          {"input": 0.80,  "output": 4.00},
    "claude-haiku-4-5-20251001": {"input": 0.80,  "output": 4.00},
    "claude-sonnet-4-6":         {"input": 3.00,  "output": 15.00},
}

def calc_cost_usd(model: str, inp: int, out: int) -> float:
    p = PRICING.get(model, {"input": 3.00, "output": 15.00})
    return round((inp * p["input"] + out * p["output"]) / 1_000_000, 6)


def _calc_xal(req) -> int:
    cost = 30
    s = req.slide_count
    if s >= 6 and s <= 10: cost += 5
    elif s == 15: cost += 8
    elif s == 20: cost += 10
    elif s == 25: cost += 12
    elif s == 30: cost += 15
    if req.detail_level == "Detailed": cost += 3
    cost += _addon_count(req) * 0  # addons don't add cost in current formula
    return int(cost)

def smart_estimate(slide_count: int, detail_level: str, addon_count: int, is_pro: bool) -> int:
    """
    Smart time estimator. Pro uses claude-sonnet (slower, higher quality),
    Free uses claude-haiku (faster). Always keeps the two plans separated.
    Falls back to a calibrated formula when not enough same-plan data exists.
    """
    stats = _load_stats()

    # Filter to only same-plan entries for a cleaner signal
    same_plan = [e for e in stats if e.get("is_pro") == is_pro]

    def formula_estimate() -> int:
        # Calibrated separately per plan — sonnet is roughly 1.8-2x slower than haiku
        base = 45 if is_pro else 25
        t = base + max(0, slide_count - 5) * (3.5 if is_pro else 2.0)
        if detail_level == "Detailed": t += 18 if is_pro else 10
        if detail_level == "Summary":  t -= 5
        t += addon_count * (5 if is_pro else 3)
        return max(int(t), 20)

    # Need at least 5 same-plan points for KNN to be meaningful
    if len(same_plan) < 5:
        return formula_estimate()

    # KNN on same-plan data only — no cross-plan contamination
    d_score = _detail_score(detail_level)
    scored = []
    for i, entry in enumerate(same_plan):
        dist = (
            abs(entry["slide_count"] - slide_count) * 3.0 +
            abs(_detail_score(entry["detail_level"]) - d_score) * 5.0 +
            abs(entry["addon_count"] - addon_count) * 2.0
        )
        # Newer entries weighted higher (index grows toward end = newer)
        recency_weight = 1.0 + (i / len(same_plan)) * 0.6
        scored.append((dist, recency_weight, entry["actual_seconds"]))

    scored.sort(key=lambda x: x[0])
    # Use top 7 neighbors for stability (more data = smoother estimate)
    top = scored[:7]
    total_w, total_wt = 0.0, 0.0
    for dist, rw, secs in top:
        w = rw / (dist + 1.0)
        total_wt += w * secs
        total_w += w
    knn_estimate = int(total_wt / total_w) if total_w > 0 else formula_estimate()

    # Blend KNN with formula when data is limited (< 20 same-plan points)
    if len(same_plan) < 20:
        blend = len(same_plan) / 20.0          # 0.25 … 0.99
        knn_estimate = int(knn_estimate * blend + formula_estimate() * (1 - blend))

    return max(knn_estimate, 15)

def record_generation(slide_count: int, detail_level: str, addon_count: int, is_pro: bool, actual_seconds: float):
    stats = _load_stats()
    stats.append({
        "slide_count": slide_count,
        "detail_level": detail_level,
        "addon_count": addon_count,
        "is_pro": is_pro,
        "actual_seconds": round(actual_seconds),
        "ts": int(time.time()),
    })
    _save_stats(stats)

# ─────────────────────────────────────────────────────────────────────────────


@app.get("/admin/users")
async def admin_users(secret: str = ""):
    if secret != os.getenv("ADMIN_SECRET", "raportakam-admin-2026"):
        raise HTTPException(status_code=403, detail="Forbidden")
    try:
        from supabase import create_client
        sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))

        # Run both queries in parallel via threads
        profiles_res, gen_res = await asyncio.gather(
            asyncio.to_thread(lambda: sb.table("profiles").select("*").execute()),
            asyncio.to_thread(lambda: sb.rpc("get_user_gen_stats", {}).execute()),
        )
        profiles = {p["id"]: p for p in profiles_res.data}

        # gen_res.data = [{ user_id, total_tokens, total_gens }]
        token_map = {r["user_id"]: r["total_tokens"] for r in (gen_res.data or [])}
        gen_count_map = {r["user_id"]: r["total_gens"] for r in (gen_res.data or [])}
        cost_map = {r["user_id"]: float(r.get("total_cost", 0) or 0) for r in (gen_res.data or [])}

        from datetime import datetime, timezone
        now = datetime.now(timezone.utc)
        users = []
        for uid, p in profiles.items():
            plan = p.get("plan", "free")
            expires_at = p.get("plan_expires_at")
            if plan == "pro" and expires_at:
                exp = datetime.fromisoformat(expires_at.replace("Z", "+00:00"))
                if exp < now:
                    plan = "free"
                    sb.table("profiles").update({"plan": "free", "plan_expires_at": None}).eq("id", uid).execute()
                    expires_at = None
            users.append({
                "id": uid,
                "email": p.get("email", ""),
                "full_name": p.get("full_name", ""),
                "plan": plan,
                "plan_expires_at": expires_at,
                "points": p.get("points", 100),
                "generations_used": gen_count_map.get(uid, 0),
                "tokens_used": token_map.get(uid, 0),
                "cost_usd": cost_map.get(uid, 0.0),
                "created_at": p.get("created_at", ""),
                "user_number": p.get("user_number"),
            })

        users.sort(key=lambda x: x["tokens_used"], reverse=True)
        return {"users": users, "total": len(users)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/admin/set-plan")
async def set_plan(secret: str = "", user_id: str = "", plan: str = ""):
    if secret != os.getenv("ADMIN_SECRET", "raportakam-admin-2026"):
        raise HTTPException(status_code=403, detail="Forbidden")
    if plan not in ("free", "pro"):
        raise HTTPException(status_code=400, detail="Invalid plan")
    try:
        from supabase import create_client
        from datetime import datetime, timezone, timedelta
        sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))
        expires_at = (datetime.now(timezone.utc) + timedelta(days=30)).isoformat() if plan == "pro" else None
        sb.table("profiles").update({"plan": plan, "plan_expires_at": expires_at}).eq("id", user_id).execute()
        return {"success": True, "plan_expires_at": expires_at}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generation/delete")
async def delete_generation(generation_id: str = "", user_id: str = ""):
    """Soft-delete a generation — keeps row for stats, hides from user history."""
    if not generation_id or not user_id:
        raise HTTPException(status_code=400, detail="Missing generation_id or user_id")
    try:
        from supabase import create_client
        sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))
        sb.table("generations").update({"deleted": True}).eq("id", generation_id).eq("user_id", user_id).execute()
        return {"success": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/admin/add-points")
async def add_points(secret: str = "", user_id: str = "", points: int = 0):
    if secret != os.getenv("ADMIN_SECRET", "raportakam-admin-2026"):
        raise HTTPException(status_code=403, detail="Forbidden")
    try:
        from supabase import create_client
        sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))
        profile = sb.table("profiles").select("points").eq("id", user_id).single().execute()
        current = profile.data.get("points", 100) or 0
        new_points = current + points
        sb.table("profiles").update({"points": new_points}).eq("id", user_id).execute()
        return {"success": True, "points": new_points}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def fix_text_overflow(path):
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    tf.word_wrap = True
                    total_chars = sum(len(p.text) for p in tf.paragraphs)
                    if total_chars > 80:
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        prs.save(path)
    except Exception as e:
        print(f"[fix_text_overflow] skipped: {e}")


def sanitize_js(code: str) -> str:
    # Fix wrong pptxgenjs constructor variants → new pptxgen()
    code = re.sub(r'new\s+pptxgen\.PresentationOptions\s*\(\s*\)', 'new pptxgen()', code)
    code = re.sub(r'new\s+pptxgen\.Presentation\s*\(\s*\)', 'new pptxgen()', code)
    code = re.sub(r'new\s+pptxgen\.PptxGenJS\s*\(\s*\)', 'new pptxgen()', code)
    code = re.sub(r'new\s+PptxGenJS\s*\(\s*\)', 'new pptxgen()', code)
    # Curly/smart quotes → straight ASCII
    code = code.replace('\u201c', '"').replace('\u201d', '"')
    code = code.replace('\u2018', "'").replace('\u2019', "'")
    # Remove # from hex color strings: "#RRGGBB" → "RRGGBB"
    code = re.sub(r'(["\'])#([0-9a-fA-F]{3,6})\1', r'\1\2\1', code)
    # Strip 8-char hex colors down to 6 (opacity-in-hex is invalid)
    code = re.sub(r'(["\'])([0-9a-fA-F]{8})\1', lambda m: f'{m.group(1)}{m.group(2)[:6]}{m.group(1)}', code)
    # Fix valig → valign
    code = re.sub(r'\bvalig\b(?!n)', 'valign', code)
    # Fix align: "justified" → align: "left" (not supported by pptxgenjs)
    code = re.sub(r'align:\s*["\']justified["\']', 'align: "left"', code)
    # Replace LINE shape type with rect to avoid missing-param errors
    code = re.sub(r'pptxgen\.shapes\.LINE', 'pres.ShapeType.rect', code, flags=re.IGNORECASE)
    code = re.sub(r'pres\.ShapeType\.line\b', 'pres.ShapeType.rect', code, flags=re.IGNORECASE)
    code = re.sub(r'ShapeType\.line\b', 'pres.ShapeType.rect', code, flags=re.IGNORECASE)
    # Fix h:0 (bare zero height) → h:0.02 so rect renders, but don't touch h:0.07 etc.
    code = re.sub(r'\bh:\s*0(?![\.\d])', 'h:0.02', code)
    # Replace backtick template literals with empty string fallback (rare but happens)
    code = re.sub(r'`[^`]*`', lambda m: '"' + m.group(0)[1:-1].replace('"', '\\"').replace('${', '"+').replace('}', '+"') + '"', code)
    # Fix addChart calls with empty data arrays to avoid PptxGenJS crash
    code = re.sub(
        r'(\.addChart\s*\([^,]+,\s*)\[\s*\]',
        r'\1[{name:"Data",labels:["A","B","C","D"],values:[10,20,30,40]}]',
        code
    )
    return code


def pick_model_and_tokens(slide_count: int, is_pro: bool, addon_count: int = 0):
    model = "claude-sonnet-4-6" if is_pro else "claude-haiku-4-5-20251001"
    if slide_count <= 5:    tokens = 10000
    elif slide_count <= 10: tokens = 13000
    elif slide_count <= 15: tokens = 16000
    elif slide_count <= 20: tokens = 19000
    elif slide_count <= 25: tokens = 22000
    else:                   tokens = 26000
    # Each addon adds extra slides/content — bump the budget to avoid truncation
    tokens += addon_count * 600
    return model, tokens


class GenerateRequest(BaseModel):
    topic: str
    slide_count: int = 10
    color_theme: str = "Auto (AI decides)"
    language: str = "English"
    style: str = "Academic / University"
    detail_level: str = "Balanced"
    student_name: str = ""
    instructor_name: str = ""
    university_name: str = ""
    date: str = "2026"
    addon_table: bool = False
    addon_timeline: bool = False
    addon_chart_extra: bool = False
    addon_quotes: bool = False
    addon_comparison: bool = False
    addon_cover_page: bool = False
    addon_custom_text: str = ""
    conclusion: bool = True
    addon_references: bool = False
    file_name: str = "raportakam"
    is_pro: bool = False
    user_id: str = ""
    plan: str = "free"


def build_prompt(req: GenerateRequest) -> str:
    prompt = SYSTEM_PROMPT
    replacements = {
        "{{TOPIC}}":             req.topic,
        "{{SLIDE_COUNT}}":       str(req.slide_count),
        "{{COLOR_THEME}}":       req.color_theme,
        "{{LANGUAGE}}":          req.language,
        "{{STYLE}}":             req.style,
        "{{DETAIL_LEVEL}}":      req.detail_level,
        "{{STUDENT_NAME}}":      req.student_name,
        "{{INSTRUCTOR_NAME}}":   req.instructor_name,
        "{{UNIVERSITY_NAME}}":   req.university_name,
        "{{DATE}}":              req.date,
        "{{ADDON_TABLE}}":       "ON" if req.addon_table else "OFF",
        "{{ADDON_TIMELINE}}":    "ON" if req.addon_timeline else "OFF",
        "{{ADDON_CHART_EXTRA}}": "ON" if req.addon_chart_extra else "OFF",
        "{{ADDON_QUOTES}}":      "ON" if req.addon_quotes else "OFF",
        "{{ADDON_COMPARISON}}":  "ON" if req.addon_comparison else "OFF",
        "{{ADDON_COVER_PAGE}}":  "ON" if req.addon_cover_page else "OFF",
        "{{ADDON_CUSTOM_TEXT}}": req.addon_custom_text,
        "{{CONCLUSION_SLIDE}}":   "on" if req.conclusion else "off",
        "{{ADDON_CONCLUSION}}":   "ON" if req.conclusion else "OFF",
        "{{ADDON_REFERENCES}}":   "ON" if req.addon_references else "OFF",
    }
    for k, v in replacements.items():
        prompt = prompt.replace(k, v)
    return prompt


r2 = boto3.client(
    "s3",
    endpoint_url=os.getenv("R2_ENDPOINT"),
    aws_access_key_id=os.getenv("R2_ACCESS_KEY_ID"),
    aws_secret_access_key=os.getenv("R2_SECRET_ACCESS_KEY"),
    config=Config(signature_version="s3v4"),
    region_name="auto",
)
BUCKET = os.getenv("R2_BUCKET_NAME")


@app.get("/estimate")
async def estimate_time(slide_count: int = 10, detail_level: str = "Balanced", addon_count: int = 0, is_pro: bool = False):
    secs = smart_estimate(slide_count, detail_level, addon_count, is_pro)
    return {"seconds": secs}


@app.post("/generate/stream")
async def generate_stream(req: GenerateRequest):
    prompt = build_prompt(req)
    _gen_start = time.time()

    async def event_stream():
        yield f"data: {json.dumps({'stage': 1, 'label': 'Analyzing topic...'})}\n\n"
        full_code = ""
        char_count = 0
        tokens_used = 0
        input_tokens = 0
        output_tokens = 0
        _model_used, _ = pick_model_and_tokens(req.slide_count, req.is_pro, _addon_count(req))
        try:
            def run_claude():
                nonlocal full_code, char_count, tokens_used, input_tokens, output_tokens
                _model, _tokens = pick_model_and_tokens(req.slide_count, req.is_pro, _addon_count(req))
                with claude.messages.stream(
                    model=_model,
                    max_tokens=_tokens,
                    messages=[{"role": "user", "content": prompt}]
                ) as stream:
                    for text in stream.text_stream:
                        full_code += text
                        char_count += len(text)
                    usage = stream.get_final_message().usage
                    input_tokens = usage.input_tokens or 0
                    output_tokens = usage.output_tokens or 0
                    tokens_used = input_tokens + output_tokens
            await asyncio.to_thread(run_claude)
            print(f"[Claude] chars={char_count} tokens={tokens_used} has_writefile={'pres.writeFile' in full_code} last50={repr(full_code[-50:])}")
        except Exception as e:
            yield f"data: {json.dumps({'stage': -1, 'error': str(e)})}\n\n"
            return

        yield f"data: {json.dumps({'stage': 2, 'label': 'Designing layout...'})}\n\n"

        def _extract_code(raw: str) -> str:
            raw = raw.strip()
            raw = re.sub(r"^```[a-z]*\n?", "", raw)
            raw = re.sub(r"\n?```$", "", raw).strip()
            js_start = raw.find('const pptxgen = require')
            if js_start == -1:
                js_start = raw.find('const pptxgen=require')
            if js_start > 0:
                raw = raw[js_start:]
            return sanitize_js(raw.strip())

        clean_code = _extract_code(full_code)

        # Detect truncation: code must end with writeFile call
        if 'pres.writeFile' not in clean_code:
            print(f"[TRUNCATED] last100={repr(clean_code[-100:])}")
            yield f"data: {json.dumps({'stage': -1, 'error': 'کۆدەکە ناتەواو بوو. دووبارە هەوڵ بدەرەوە.'})}\n\n"
            return

        yield f"data: {json.dumps({'stage': 3, 'label': 'Generating slides...'})}\n\n"

        node_modules = os.path.join(os.path.dirname(__file__), "node_modules")
        backend_dir = os.path.dirname(__file__).replace(chr(92), "/")

        async def run_node(code_to_run: str):
            with tempfile.TemporaryDirectory() as tmpdir:
                js_path = os.path.join(tmpdir, "gen.js")
                out_path = os.path.join(tmpdir, "output.pptx")
                out_path_js = out_path.replace('\\', '/')
                code = re.sub(
                    r'fileName\s*:\s*["\']output\.pptx["\']',
                    f'fileName: "{out_path_js}"',
                    code_to_run
                )
                with open(js_path, "w", encoding="utf-8") as f:
                    f.write(f'process.chdir("{backend_dir}"); {code}')
                js_path_copy = js_path  # avoid closure issues
                def _run():
                    return subprocess.run(
                        ["node", js_path_copy],
                        capture_output=True, text=True, timeout=90,
                        env={**os.environ, "NODE_PATH": node_modules}
                    )
                res = await asyncio.to_thread(_run)
                if res.returncode == 0:
                    saved = os.path.join(tempfile.gettempdir(), f"pptx_{uuid.uuid4()}.pptx")
                    shutil.copy2(out_path, saved)
                    return res, saved
                return res, None

        async def generate_code():
            nonlocal tokens_used, input_tokens, output_tokens
            code = ""
            _model, _tokens = pick_model_and_tokens(req.slide_count, req.is_pro, _addon_count(req))
            def _call():
                nonlocal code, tokens_used, input_tokens, output_tokens
                with claude.messages.stream(
                    model=_model,
                    max_tokens=_tokens,
                    messages=[{"role": "user", "content": prompt}]
                ) as stream:
                    for text in stream.text_stream:
                        code += text
                    usage = stream.get_final_message().usage
                    input_tokens += usage.input_tokens or 0
                    output_tokens += usage.output_tokens or 0
                    tokens_used += (usage.input_tokens or 0) + (usage.output_tokens or 0)
            await asyncio.to_thread(_call)
            return _extract_code(code)

        try:
            result, saved_path = await run_node(clean_code)

            # Auto-retry once if Node failed
            if result.returncode != 0:
                yield f"data: {json.dumps({'stage': 2, 'label': 'Retrying...'})}\n\n"
                try:
                    retry_code = await generate_code()
                except Exception as e:
                    yield f"data: {json.dumps({'stage': -1, 'error': str(e)})}\n\n"
                    return
                if 'pres.writeFile' not in retry_code:
                    yield f"data: {json.dumps({'stage': -1, 'error': 'کۆدەکە ناتەواو بوو. دووبارە هەوڵ بدەرەوە.'})}\n\n"
                    return
                yield f"data: {json.dumps({'stage': 3, 'label': 'Generating slides...'})}\n\n"
                result, saved_path = await run_node(retry_code)

            if result.returncode != 0:
                err = result.stderr or result.stdout or "Node.js error"
                err = re.sub(r'[A-Za-z]:[^\n:]+\.js:\d+\s*', '', err).strip()
                err = re.sub(r'\s+', ' ', err).strip()
                err = err[:200] if err else "کێشەیەک روویدا لە دروستکردنی فایلەکە"
                yield f"data: {json.dumps({'stage': -1, 'error': err})}\n\n"
                return

            yield f"data: {json.dumps({'stage': 4, 'label': 'Building your file...'})}\n\n"

            fix_text_overflow(saved_path)

            file_id = str(uuid.uuid4())
            r2_key = f"presentations/{file_id}.pptx"
            r2.upload_file(
                saved_path, BUCKET, r2_key,
                ExtraArgs={"ContentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
            )
            try:
                os.remove(saved_path)
            except Exception:
                pass
            actual_secs = time.time() - _gen_start
            record_generation(req.slide_count, req.detail_level, _addon_count(req), req.is_pro, actual_secs)
            download_url = f"/pptx/{file_id}"

            # Record generation + deduct points server-side (service role bypasses RLS)
            if req.user_id:
                try:
                    from supabase import create_client
                    sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))
                    cost_usd = calc_cost_usd(_model_used, input_tokens, output_tokens)
                    sb.table("generations").insert({
                        "user_id": req.user_id,
                        "prompt": req.topic,
                        "output_type": "pptx",
                        "status": "done",
                        "file_name": req.file_name,
                        "tokens_used": tokens_used,
                        "input_tokens": input_tokens,
                        "output_tokens": output_tokens,
                        "model": _model_used,
                        "cost_usd": cost_usd,
                        "file_url": download_url,
                    }).execute()
                    if req.plan != "pro":
                        profile = sb.table("profiles").select("points").eq("id", req.user_id).single().execute()
                        current_points = (profile.data or {}).get("points", 100)
                        sb.table("profiles").update({"points": max(0, current_points - _calc_xal(req))}).eq("id", req.user_id).execute()
                except Exception as db_err:
                    print(f"[DB] Error recording generation: {db_err}")

            yield f"data: {json.dumps({'stage': 5, 'label': 'Ready!', 'url': download_url, 'file_id': file_id, 'tokens_used': tokens_used})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'stage': -1, 'error': str(e)})}\n\n"

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.get("/pptx/{file_id}")
async def download_pptx(file_id: str, name: str = "raportakam"):
    try:
        # unique path per request to avoid conflicts
        local_path = os.path.join(tempfile.gettempdir(), f"dl_{uuid.uuid4()}.pptx")
        r2.download_file(BUCKET, f"presentations/{file_id}.pptx", local_path)
        from fastapi.responses import FileResponse
        download_name = name if name.endswith(".pptx") else f"{name}.pptx"
        return FileResponse(local_path, filename=download_name,
                            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            background=None)
    except Exception:
        raise HTTPException(status_code=404, detail="File not found")


@app.get("/health")
async def health():
    return {"status": "ok"}


DETECT_SYSTEM_PROMPT = """You are an expert AI text detection system trained to analyze writing patterns and identify whether text was written by a human or generated by an AI model (such as ChatGPT, Claude, Gemini, Copilot, etc.).

Analyze the given text and return ONLY a valid JSON object with no extra text, no markdown, no explanation outside the JSON.

Detect based on these signals:
AI indicators:
- Perfectly balanced sentence length variation
- Overuse of transitional phrases ("Furthermore", "Moreover", "In conclusion", "It is worth noting")
- Hedging language ("It is important to", "One must consider")
- Lack of personal opinions, emotions, or lived experience
- Unnaturally consistent tone throughout
- Generic examples instead of specific or personal ones
- Repetitive structure across paragraphs
- No typos, no informal contractions, no colloquialisms
- Overly formal vocabulary for the context
- Lists and enumerations used excessively
- Conclusions that summarize without adding new insight
- Missing cultural/regional flavor or slang
- Suspiciously perfect grammar and punctuation

Human indicators:
- Inconsistent tone or style shifts
- Personal anecdotes or emotional reactions
- Typos, informal language, contractions
- Tangents or off-topic thoughts
- Strong or biased opinions
- Specific names, places, dates from personal experience
- Irregular sentence rhythm
- Humor, sarcasm, or irony
- Regional slang or dialect

Return this exact JSON format:
{
  "ai_percentage": 85,
  "human_percentage": 15,
  "verdict": "Very likely AI-generated",
  "confidence": "High",
  "source_type": "plain_text",
  "top_reasons": [
    "Excessive use of transitional phrases",
    "Unnaturally consistent tone",
    "No personal voice or emotion"
  ],
  "risky_sentences": [
    "Sentence from the text that looks most AI-generated"
  ],
  "suggestion": "To make this more human, add personal examples, vary sentence length, and include informal language."
}

source_type should be one of: "plain_text", "pdf", "docx", "pptx"
Verdict scale:
- 0–20% → Almost certainly human
- 21–40% → Likely human with some AI assistance
- 41–60% → Mixed or uncertain
- 61–80% → Likely AI-generated
- 81–100% → Almost certainly AI-generated

Be strict and critical. Do not give benefit of the doubt. If it reads clean, structured, and formal — lean toward AI. Short texts under 50 words should return a low confidence score."""


@app.post("/ai-detect")
async def detect_ai(
    text: str = Form(None),
    file: UploadFile = File(None),
    use_sonnet: str = Form("0"),
    user_id: str = Form(""),
):
    import base64, io
    model = "claude-sonnet-4-6" if use_sonnet == "1" else "claude-haiku-4-5-20251001"
    client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    messages = []
    source_type = "plain_text"

    if file and file.filename:
        content = await file.read()
        fname = file.filename.lower()

        if fname.endswith(".pdf"):
            source_type = "pdf"
            pdf_b64 = base64.standard_b64encode(content).decode("utf-8")
            messages = [{
                "role": "user",
                "content": [
                    {
                        "type": "document",
                        "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64},
                    },
                    {"type": "text", "text": f"Analyze this document for AI detection. Set source_type to \"{source_type}\"."},
                ],
            }]
        elif fname.endswith(".docx"):
            source_type = "docx"
            from docx import Document as DocxDocument
            doc = DocxDocument(io.BytesIO(content))
            extracted = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            messages = [{"role": "user", "content": f"[source_type: docx]\n\n{extracted}"}]
        elif fname.endswith(".pptx"):
            source_type = "pptx"
            prs = Presentation(io.BytesIO(content))
            extracted = "\n".join([
                shape.text for slide in prs.slides
                for shape in slide.shapes if shape.has_text_frame and shape.text.strip()
            ])
            messages = [{"role": "user", "content": f"[source_type: pptx]\n\n{extracted}"}]
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF, DOCX, or PPTX.")
    elif text and text.strip():
        messages = [{"role": "user", "content": f"[source_type: plain_text]\n\n{text}"}]
    else:
        raise HTTPException(status_code=400, detail="Provide text or a file.")

    response = client.messages.create(
        model=model,
        max_tokens=1024,
        system=DETECT_SYSTEM_PROMPT,
        messages=messages,
    )
    raw = response.content[0].text.strip()
    # Strip markdown code fences if present
    if raw.startswith("```"):
        raw = re.sub(r"^```[a-z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)
    result = json.loads(raw)
    result["source_type"] = source_type

    # Record usage to DB
    if user_id:
        try:
            from supabase import create_client
            sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))
            inp = response.usage.input_tokens or 0
            out = response.usage.output_tokens or 0
            cost_usd = calc_cost_usd(model, inp, out)
            sb.table("generations").insert({
                "user_id": user_id,
                "prompt": (text or "")[:200] or (file.filename if file else "file"),
                "output_type": "ai-detect",
                "status": "done",
                "file_name": file.filename if file else None,
                "tokens_used": inp + out,
                "input_tokens": inp,
                "output_tokens": out,
                "model": model,
                "cost_usd": cost_usd,
            }).execute()
        except Exception as db_err:
            print(f"[DB] ai-detect record error: {db_err}")

    return result


# ── Report Generator ─────────────────────────────────────────────────────────

REPORT_SYSTEM = """You are an expert academic writer. Generate complete, polished, university-level academic reports with zero errors.

STRICT OUTPUT RULES:
- Output ONLY the final report — no explanations, no notes, no meta-commentary
- Use markdown: # for the main title, ## for major sections, ### for subsections, **bold** for key terms
- Never switch language mid-report
- If a field is not provided, omit it silently — do not write "N/A"

WRITING STANDARDS:
- Formal academic language throughout
- Third person only (no "I", "we", "you")
- No contractions (don't → do not, can't → cannot)
- Logical paragraph structure with proper transitions (Furthermore, However, Therefore, In contrast)
- Every paragraph covers one focused idea
- Vary sentence length and structure — avoid robotic, repetitive phrasing
- Zero grammar or spelling errors

STRUCTURE (in this order):
1. # Main Title — large title, followed by cover details (student, course, instructor, date) as plain lines
2. ## Abstract — only if requested
3. ## Introduction — context, objectives, scope
4. ## [Body Sections] — expand each provided key point into a full detailed section
5. ## Discussion — analysis and implications (skip for Short length)
6. ## Conclusion — only if Include Conclusion is Yes
7. ## References — only if Include References is Yes, formatted in the specified citation style

LENGTH TARGETS (word count for body text, excluding cover):
- Short: 500–800 words
- Medium: 1,000–1,500 words
- Long: 2,000–3,000 words

CITATION RULES:
- Basic research level: NO citations, no reference list
- Medium/Advanced: realistic in-text citations + properly formatted reference list"""


class ReportRequest(BaseModel):
    topic: str = ""
    title: str = ""
    student_name: str = ""
    course: str = ""
    instructor: str = ""
    date: str = ""
    purpose: str = ""
    points: str = ""
    length: str = "Medium"
    style: str = "Formal Academic"
    research_level: str = "Medium"
    citation_styles: list = ["APA"]
    include_abstract: bool = False
    include_references: bool = True
    include_conclusion: bool = True
    language: str = "English"
    user_id: str = ""
    plan: str = "free"


def build_report_prompt(req: ReportRequest) -> str:
    lines = ["Write a complete academic report with the following specifications:\n"]
    lines.append(f"Topic: {req.topic}")
    if req.title:        lines.append(f"Title: {req.title}")
    if req.student_name: lines.append(f"Student Name: {req.student_name}")
    if req.course:       lines.append(f"University/Course: {req.course}")
    if req.instructor:   lines.append(f"Instructor: {req.instructor}")
    if req.date:         lines.append(f"Date: {req.date}")
    if req.purpose:      lines.append(f"Purpose: {req.purpose}")
    if req.points:       lines.append(f"Key Points/Sections:\n{req.points}")
    lines.append(f"\nLength: {req.length}")
    lines.append(f"Writing Style: {req.style}")
    lines.append(f"Research Level: {req.research_level}")
    if req.research_level != "Basic" and req.citation_styles:
        lines.append(f"Citation Style: {', '.join(req.citation_styles)}")
    lines.append(f"Include Abstract: {'Yes' if req.include_abstract else 'No'}")
    lines.append(f"Include Conclusion: {'Yes' if req.include_conclusion else 'No'}")
    lines.append(f"Include References: {'Yes' if req.include_references else 'No'}")
    lines.append(f"Language: {req.language}")
    lines.append("\nGenerate the complete report now.")
    return "\n".join(lines)


def estimate_report_seconds(length: str, is_pro: bool) -> int:
    base = {"Short": 35, "Medium": 65, "Long": 120}.get(length, 65)
    return int(base * 1.8) if is_pro else base


@app.get("/report/estimate")
async def report_estimate(length: str = "Medium", is_pro: bool = False):
    return {"seconds": estimate_report_seconds(length, is_pro)}


@app.post("/report/stream")
async def report_stream(req: ReportRequest):
    if not req.topic.strip():
        raise HTTPException(status_code=400, detail="Topic is required")
    prompt = build_report_prompt(req)
    is_pro = req.plan == "pro"
    is_rtl_lang = req.language in ["Kurdish (Sorani)", "Arabic"]
    # Kurdish & Arabic always use Sonnet for quality; pro users always get Sonnet
    model = "claude-sonnet-4-6" if (is_pro or is_rtl_lang) else "claude-haiku-4-5-20251001"
    max_tokens = {"Short": 2000, "Medium": 5000, "Long": 10000}.get(req.length, 5000)

    async def event_stream():
        yield f"data: {json.dumps({'stage': 1, 'label': 'شیکردنەوەی بابەت...'})}\n\n"

        queue: asyncio.Queue = asyncio.Queue()
        loop = asyncio.get_running_loop()
        inp_tok = [0]
        out_tok = [0]
        first_chunk_sent = [False]

        def run_claude():
            try:
                with claude.messages.stream(
                    model=model,
                    max_tokens=max_tokens,
                    temperature=0.3,
                    system=REPORT_SYSTEM,
                    messages=[{"role": "user", "content": prompt}]
                ) as stream:
                    for text in stream.text_stream:
                        if not first_chunk_sent[0]:
                            asyncio.run_coroutine_threadsafe(queue.put(("stage2", None)), loop)
                            first_chunk_sent[0] = True
                        asyncio.run_coroutine_threadsafe(queue.put(("chunk", text)), loop)
                    usage = stream.get_final_message().usage
                    inp_tok[0] = usage.input_tokens or 0
                    out_tok[0] = usage.output_tokens or 0
                    asyncio.run_coroutine_threadsafe(queue.put(("done", None)), loop)
            except Exception as e:
                asyncio.run_coroutine_threadsafe(queue.put(("error", str(e))), loop)

        threading.Thread(target=run_claude, daemon=True).start()

        while True:
            kind, data = await queue.get()
            if kind == "stage2":
                yield f"data: {json.dumps({'stage': 2, 'label': 'نووسینی ناوەڕۆک...'})}\n\n"
            elif kind == "chunk":
                yield f"data: {json.dumps({'chunk': data})}\n\n"
            elif kind == "done":
                yield f"data: {json.dumps({'stage': 3, 'label': 'ئامادەکردنی ڕاپۆرت...'})}\n\n"
                tokens = inp_tok[0] + out_tok[0]
                cost = calc_cost_usd(model, inp_tok[0], out_tok[0])
                if req.user_id:
                    try:
                        from supabase import create_client
                        sb = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_SECRET_KEY"))
                        sb.table("generations").insert({
                            "user_id": req.user_id,
                            "prompt": req.topic[:200],
                            "output_type": "report",
                            "status": "done",
                            "file_name": req.title or req.topic[:50],
                            "tokens_used": tokens,
                            "input_tokens": inp_tok[0],
                            "output_tokens": out_tok[0],
                            "model": model,
                            "cost_usd": cost,
                        }).execute()
                        if req.plan != "pro":
                            report_cost = (30
                                + (2 if req.length == "Long" else 0)
                                + (2 if req.style == "Advanced Academic" else 0)
                                + (2 if req.research_level == "Advanced" else 0)
                                + (2 if req.include_abstract else 0)
                                + (2 if req.include_references else 0))
                            profile = sb.table("profiles").select("points").eq("id", req.user_id).single().execute()
                            current = (profile.data or {}).get("points", 100)
                            sb.table("profiles").update({"points": max(0, current - report_cost)}).eq("id", req.user_id).execute()
                    except Exception as db_err:
                        print(f"[DB] report record error: {db_err}")
                yield f"data: {json.dumps({'done': True, 'tokens': tokens})}\n\n"
                break
            elif kind == "error":
                yield f"data: {json.dumps({'error': data})}\n\n"
                break

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


import random as _random

KURDISH_FONTS = ["Rabar_022", "NRT"]

# Layouts differ in: font, size, spacing, margins, heading alignment, underline, caps, indent, border
WORD_TEMPLATES = [
    {   # Classic academic — centered title with rule, underlined h2, generous spacing
        "body_font": "Times New Roman", "heading_font": "Times New Roman",
        "h1_size": 30, "h2_size": 18, "h3_size": 14, "body_size": 13,
        "h1_bold": True,  "h2_bold": True,  "h3_bold": True,
        "h1_italic": False, "h2_italic": False, "h3_italic": False,
        "h1_center": True,  "h2_center": False,
        "h2_underline": True, "h2_caps": False,
        "h2_border": False, "h1_border": True,
        "line_spacing_pt": 22, "space_after_pt": 10,
        "margin_in": 1.25, "first_indent": False,
    },
    {   # Modern minimal — Calibri, left-aligned, underline-less, tight
        "body_font": "Calibri", "heading_font": "Calibri",
        "h1_size": 32, "h2_size": 17, "h3_size": 13, "body_size": 12,
        "h1_bold": True,  "h2_bold": True,  "h3_bold": False,
        "h1_italic": False, "h2_italic": False, "h3_italic": False,
        "h1_center": True,  "h2_center": False,
        "h2_underline": False, "h2_caps": False,
        "h2_border": True, "h1_border": False,
        "line_spacing_pt": 16, "space_after_pt": 7,
        "margin_in": 1.0, "first_indent": False,
    },
    {   # Elegant literary — Georgia, wide margins, double spacing, italic h3, indent
        "body_font": "Georgia", "heading_font": "Georgia",
        "h1_size": 30, "h2_size": 18, "h3_size": 14, "body_size": 13,
        "h1_bold": True,  "h2_bold": True,  "h3_bold": False,
        "h1_italic": False, "h2_italic": False, "h3_italic": True,
        "h1_center": True,  "h2_center": True,
        "h2_underline": False, "h2_caps": False,
        "h2_border": False, "h1_border": True,
        "line_spacing_pt": 26, "space_after_pt": 12,
        "margin_in": 1.5, "first_indent": True,
    },
    {   # Technical — Arial, ALL CAPS h2, compact, narrow margins
        "body_font": "Arial", "heading_font": "Arial",
        "h1_size": 24, "h2_size": 15, "h3_size": 12, "body_size": 11,
        "h1_bold": True,  "h2_bold": True,  "h3_bold": True,
        "h1_italic": False, "h2_italic": False, "h3_italic": False,
        "h1_center": False, "h2_center": False,
        "h2_underline": False, "h2_caps": True,
        "h2_border": False, "h1_border": False,
        "line_spacing_pt": 14, "space_after_pt": 6,
        "margin_in": 1.0, "first_indent": False,
    },
    {   # Airy open — Calibri, very loose, h2 non-bold with border, huge title
        "body_font": "Calibri", "heading_font": "Calibri",
        "h1_size": 36, "h2_size": 16, "h3_size": 13, "body_size": 12,
        "h1_bold": True,  "h2_bold": False, "h3_bold": False,
        "h1_italic": False, "h2_italic": False, "h3_italic": False,
        "h1_center": True,  "h2_center": False,
        "h2_underline": False, "h2_caps": False,
        "h2_border": True, "h1_border": False,
        "line_spacing_pt": 24, "space_after_pt": 12,
        "margin_in": 1.5, "first_indent": False,
    },
]


@app.post("/report/download/word")
async def report_download_word(
    text: str = Form(...),
    filename: str = Form("report"),
    language: str = Form("English"),
):
    import io
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    is_rtl = any(lang in language.lower() for lang in ["kurdish", "arabic", "sorani", "persian", "فارسی", "عربی", "کوردی"])
    tpl = _random.choice(WORD_TEMPLATES)
    # Kurdish/Arabic: override fonts with Rabar or NRT
    kfont = _random.choice(KURDISH_FONTS) if is_rtl else None

    def _font(key):
        return kfont if is_rtl else tpl[key]

    doc = Document()
    m = Inches(tpl["margin_in"])
    for section in doc.sections:
        section.top_margin = Inches(1.1)
        section.bottom_margin = Inches(1.1)
        section.left_margin = m
        section.right_margin = m

    BLACK = RGBColor(0, 0, 0)

    def set_rtl_para(para):
        if not is_rtl:
            return
        pPr = para._p.get_or_add_pPr()
        bidi = OxmlElement("w:bidi")
        pPr.append(bidi)
        para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    def apply_font(run, size, bold=False, italic=False, underline=False, caps=False, font_name=None):
        run.font.name = font_name or _font("body_font")
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        run.font.all_caps = caps
        run.font.color.rgb = BLACK

    def add_bottom_border(para, thick=False):
        pPr = para._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "12" if thick else "6")
        bottom.set(qn("w:space"), "4")
        bottom.set(qn("w:color"), "000000")
        pBdr.append(bottom)
        pPr.append(pBdr)

    def add_page_break():
        p = doc.add_paragraph()
        run = p.add_run()
        br = OxmlElement("w:br")
        br.set(qn("w:type"), "page")
        run._r.append(br)

    for line in text.split("\n"):
        stripped = line.strip()

        if stripped.startswith("# "):
            p = doc.add_paragraph()
            align = WD_ALIGN_PARAGRAPH.RIGHT if is_rtl else (WD_ALIGN_PARAGRAPH.CENTER if tpl["h1_center"] else WD_ALIGN_PARAGRAPH.LEFT)
            p.alignment = align
            p.paragraph_format.space_before = Pt(80)
            p.paragraph_format.space_after = Pt(28)
            run = p.add_run(stripped[2:])
            apply_font(run, tpl["h1_size"], bold=tpl["h1_bold"], italic=tpl["h1_italic"], font_name=_font("heading_font"))
            if tpl["h1_border"]:
                add_bottom_border(p, thick=True)
            set_rtl_para(p)

        elif stripped.startswith("## "):
            add_page_break()
            p = doc.add_paragraph()
            align = WD_ALIGN_PARAGRAPH.RIGHT if is_rtl else (WD_ALIGN_PARAGRAPH.CENTER if tpl["h2_center"] else WD_ALIGN_PARAGRAPH.LEFT)
            p.alignment = align
            p.paragraph_format.space_before = Pt(10)
            p.paragraph_format.space_after = Pt(6)
            run = p.add_run(stripped[3:])
            apply_font(run, tpl["h2_size"], bold=tpl["h2_bold"], italic=tpl["h2_italic"],
                       underline=tpl["h2_underline"], caps=tpl["h2_caps"], font_name=_font("heading_font"))
            if tpl["h2_border"]:
                add_bottom_border(p)
            set_rtl_para(p)

        elif stripped.startswith("### "):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(3)
            run = p.add_run(stripped[4:])
            apply_font(run, tpl["h3_size"], bold=tpl["h3_bold"], italic=tpl["h3_italic"], font_name=_font("heading_font"))
            set_rtl_para(p)

        elif stripped == "" or stripped == "---":
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(0)

        else:
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(tpl["space_after_pt"])
            p.paragraph_format.line_spacing = Pt(tpl["line_spacing_pt"])
            if tpl["first_indent"]:
                p.paragraph_format.first_line_indent = Inches(0.3)
            set_rtl_para(p)
            parts = re.split(r"(\*\*[^*]+\*\*)", stripped)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    run = p.add_run(part[2:-2])
                    apply_font(run, tpl["body_size"], bold=True)
                else:
                    run = p.add_run(part)
                    apply_font(run, tpl["body_size"])

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    safe = (filename or "report").replace(" ", "_")
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{safe}.docx"'},
    )
