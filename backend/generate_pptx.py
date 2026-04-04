from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import os
import random

THEMES = {
    "dark":     {"bg": RGBColor(10, 10, 15),     "title": RGBColor(255, 255, 255), "text": RGBColor(200, 200, 200), "accent": RGBColor(234, 179, 8)},
    "light":    {"bg": RGBColor(250, 250, 255),  "title": RGBColor(15, 23, 42),    "text": RGBColor(51, 65, 85),    "accent": RGBColor(37, 99, 235)},
    "yellow":   {"bg": RGBColor(28, 20, 5),      "title": RGBColor(251, 191, 36),  "text": RGBColor(253, 230, 138), "accent": RGBColor(245, 158, 11)},
    "blue":     {"bg": RGBColor(10, 20, 50),     "title": RGBColor(147, 197, 253), "text": RGBColor(219, 234, 254), "accent": RGBColor(59, 130, 246)},
    "green":    {"bg": RGBColor(2, 40, 30),      "title": RGBColor(110, 231, 183), "text": RGBColor(209, 250, 229), "accent": RGBColor(16, 185, 129)},
    "purple":   {"bg": RGBColor(30, 10, 50),     "title": RGBColor(216, 180, 254), "text": RGBColor(245, 243, 255), "accent": RGBColor(168, 85, 247)},
    "red":      {"bg": RGBColor(50, 0, 0),       "title": RGBColor(252, 165, 165), "text": RGBColor(254, 226, 226), "accent": RGBColor(220, 38, 38)},
    "orange":   {"bg": RGBColor(45, 10, 20),     "title": RGBColor(253, 186, 116), "text": RGBColor(255, 228, 230), "accent": RGBColor(244, 63, 94)},
    "pink":     {"bg": RGBColor(60, 20, 40),     "title": RGBColor(251, 164, 181), "text": RGBColor(255, 241, 242), "accent": RGBColor(225, 29, 72)},
    "teal":     {"bg": RGBColor(30, 30, 35),     "title": RGBColor(226, 232, 240), "text": RGBColor(241, 245, 249), "accent": RGBColor(148, 163, 184)},
    "indigo":   {"bg": RGBColor(8, 8, 28),       "title": RGBColor(167, 139, 250), "text": RGBColor(221, 214, 254), "accent": RGBColor(139, 92, 246)},
    "navy":     {"bg": RGBColor(8, 47, 73),      "title": RGBColor(103, 232, 249), "text": RGBColor(207, 250, 254), "accent": RGBColor(6, 182, 212)},
    "gold":     {"bg": RGBColor(40, 30, 10),     "title": RGBColor(251, 191, 36),  "text": RGBColor(254, 243, 199), "accent": RGBColor(217, 119, 6)},
    "forest":   {"bg": RGBColor(5, 20, 10),      "title": RGBColor(187, 247, 208), "text": RGBColor(240, 253, 244), "accent": RGBColor(22, 101, 52)},
    "sunset":   {"bg": RGBColor(60, 45, 20),     "title": RGBColor(245, 158, 11),  "text": RGBColor(254, 243, 199), "accent": RGBColor(217, 119, 6)},
    "ocean":    {"bg": RGBColor(224, 242, 254),  "title": RGBColor(3, 105, 161),   "text": RGBColor(12, 74, 110),    "accent": RGBColor(56, 189, 248)},
    "midnight": {"bg": RGBColor(40, 10, 40),     "title": RGBColor(232, 121, 249), "text": RGBColor(250, 232, 255), "accent": RGBColor(192, 38, 211)},
    "rose":     {"bg": RGBColor(53, 5, 20),      "title": RGBColor(251, 113, 133), "text": RGBColor(253, 164, 175), "accent": RGBColor(244, 63, 94)},
    "brown":    {"bg": RGBColor(30, 15, 5),      "title": RGBColor(217, 119, 6),   "text": RGBColor(253, 230, 138), "accent": RGBColor(120, 53, 15)},
    "silver":   {"bg": RGBColor(40, 40, 45),     "title": RGBColor(212, 212, 216), "text": RGBColor(212, 212, 216), "accent": RGBColor(113, 113, 122)},
}

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(tf, text: str, size: int, bold: bool, color: RGBColor, align=PP_ALIGN.RIGHT):
    tf.text = ""
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def _add_image(slide, s, img_x, img_y, img_w, img_h):
    img_path = s.get("image_path")
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, img_x, img_y, img_w, img_h)
        except Exception:
            pass

def generate_presentation(data: dict, output_path: str = None, theme_id: str = "dark", style: str = "classic", lang: str = "kurdish") -> str:
    template_path = os.path.join(os.path.dirname(__file__), "templates", f"{theme_id}.pptx")
    use_template = os.path.exists(template_path)
    
    if use_template:
        print(f"Using Master Template: {theme_id}")
        prs = Presentation(template_path)
    else:
        theme = THEMES.get(theme_id, THEMES["dark"])
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    align = PP_ALIGN.RIGHT if lang == "kurdish" else PP_ALIGN.LEFT
    opp_align = PP_ALIGN.LEFT if lang == "kurdish" else PP_ALIGN.RIGHT

    title_text = data.get("title", "ڕاپۆرتەکام")
    subtitle_text = "دروستکراوە بە Raportakam.com"

    # ── Title slide ──────────────────────────────────────────
    if use_template:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        
        # Clear hardcoded raw text from Gamma that isn't a placeholder
        for shape in list(slide.shapes):
            if not shape.is_placeholder and hasattr(shape, "text") and shape.has_text_frame:
                if shape.text.strip():
                    sp = shape.element
                    sp.getparent().remove(sp)

        # Try to find exactly 2 text placeholders (one for title, one for subtitle)
        text_phs = [ph for ph in slide.placeholders if ph.placeholder_format.type != 18]
        if len(text_phs) >= 2:
            text_phs[0].text = title_text
            text_phs[1].text = subtitle_text
        elif len(text_phs) == 1:
            text_phs[0].text = title_text + "\n" + subtitle_text
            
        for ph in text_phs:
            if ph.has_text_frame:
                for paragraph in ph.text_frame.paragraphs:
                    paragraph.alignment = align
                    for run in paragraph.runs:
                        run.font.color.rgb = theme["title"]
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide, theme["bg"])

        if style == "corporate":
            sb_x = Inches(12.88) if lang == "kurdish" else Inches(0)
            tb_x = Inches(0.5) if lang == "kurdish" else Inches(0.8)
            sidebar = slide.shapes.add_shape(1, sb_x, Inches(0), Inches(0.45), Inches(7.5))
            sidebar.fill.solid(); sidebar.fill.fore_color.rgb = theme["accent"]; sidebar.line.fill.background()
            tb = slide.shapes.add_textbox(tb_x, Inches(1.8), Inches(11.5), Inches(2.5))
            tf = tb.text_frame; tf.word_wrap = True
            add_text(tf, title_text, 40, True, theme["title"], align=align)
            line_x = Inches(7.5) if lang == "kurdish" else Inches(0.8)
            line = slide.shapes.add_shape(1, line_x, Inches(4.5), Inches(5), Inches(0.05))
            line.fill.solid(); line.fill.fore_color.rgb = theme["accent"]; line.line.fill.background()
            tb2 = slide.shapes.add_textbox(tb_x, Inches(4.8), Inches(11.5), Inches(0.8))
            add_text(tb2.text_frame, subtitle_text, 18, False, theme["text"], align=align)

        elif style == "bold":
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(3.5))
            header.fill.solid(); header.fill.fore_color.rgb = theme["accent"]; header.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(2.8))
            tf = tb.text_frame; tf.word_wrap = True
            add_text(tf, title_text, 46, True, theme["bg"], align=PP_ALIGN.CENTER)
            tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.73), Inches(0.8))
            add_text(tb2.text_frame, subtitle_text, 20, False, theme["text"], align=PP_ALIGN.CENTER)

        elif style == "minimal":
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.33), Inches(2.5))
            tf = tb.text_frame; tf.word_wrap = True
            add_text(tf, title_text, 44, True, theme["title"], align=PP_ALIGN.CENTER)
            line = slide.shapes.add_shape(1, Inches(4), Inches(4.6), Inches(5.33), Inches(0.04))
            line.fill.solid(); line.fill.fore_color.rgb = theme["accent"]; line.line.fill.background()
            tb2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.9), Inches(10.33), Inches(0.8))
            add_text(tb2.text_frame, subtitle_text, 18, False, theme["text"], align=PP_ALIGN.CENTER)

        elif style == "tech":
            top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
            top.fill.solid(); top.fill.fore_color.rgb = theme["accent"]; top.line.fill.background()
            bot = slide.shapes.add_shape(1, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12))
            bot.fill.solid(); bot.fill.fore_color.rgb = theme["accent"]; bot.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.33), Inches(2.0))
            tf = tb.text_frame; tf.word_wrap = True
            add_text(tf, title_text, 46, True, theme["title"], align=PP_ALIGN.CENTER)
            div = slide.shapes.add_shape(1, Inches(2.5), Inches(3.4), Inches(8.33), Inches(0.05))
            div.fill.solid(); div.fill.fore_color.rgb = theme["accent"]; div.line.fill.background()
            tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11.33), Inches(0.8))
            add_text(tb2.text_frame, subtitle_text, 20, False, theme["text"], align=PP_ALIGN.CENTER)
            for cx, cy in [(Inches(0.3), Inches(0.3)), (Inches(12.7), Inches(0.3))]:
                sq = slide.shapes.add_shape(1, cx, cy, Inches(0.18), Inches(0.18))
                sq.fill.solid(); sq.fill.fore_color.rgb = theme["accent"]; sq.line.fill.background()

        elif style == "elegant":
            sb_x = Inches(0) if lang == "kurdish" else Inches(12.88)
            tb_x = Inches(0.8) if lang == "kurdish" else Inches(0.5)
            sidebar = slide.shapes.add_shape(1, sb_x, Inches(0), Inches(0.45), Inches(7.5))
            sidebar.fill.solid(); sidebar.fill.fore_color.rgb = theme["accent"]; sidebar.line.fill.background()
            tb = slide.shapes.add_textbox(tb_x, Inches(1.8), Inches(12.0), Inches(2.5))
            tf = tb.text_frame; tf.word_wrap = True
            add_text(tf, title_text, 42, True, theme["title"], align=align)
            line_x = Inches(0.5) if lang == "kurdish" else Inches(7)
            line = slide.shapes.add_shape(1, line_x, Inches(4.5), Inches(5.5), Inches(0.05))
            line.fill.solid(); line.fill.fore_color.rgb = theme["accent"]; line.line.fill.background()
            tb2 = slide.shapes.add_textbox(tb_x, Inches(4.8), Inches(12.0), Inches(0.8))
            add_text(tb2.text_frame, subtitle_text, 18, False, theme["text"], align=align)

        else:
            bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.33), Inches(0.06))
            bar.fill.solid(); bar.fill.fore_color.rgb = theme["accent"]; bar.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11.33), Inches(2.0))
            tf = tb.text_frame; tf.word_wrap = True
            add_text(tf, title_text, 44, True, theme["title"], align=PP_ALIGN.CENTER)
            tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.33), Inches(0.8))
            add_text(tb2.text_frame, subtitle_text, 20, False, theme["text"], align=PP_ALIGN.CENTER)

    # ── Content slides ────────────────────────────────────────
    slides_list = data.get("slides", [])
    
    # Pre-calculate available layouts for templates
    if use_template:
        img_layouts = []
        txt_layouts = []
        for l in prs.slide_layouts:
            has_pic = any(ph.placeholder_format.type == 18 for ph in l.placeholders)
            if has_pic:
                img_layouts.append(l)
            else:
                txt_layouts.append(l)

    for idx, s in enumerate(slides_list):
        content = s.get("content", [])
        img_keywords = s.get("image_keywords")
        has_img = bool(img_keywords)

        if use_template:
            # Smart Layout Selection
            if has_img and img_layouts:
                l_to_use = random.choice(img_layouts)
            elif txt_layouts:
                # exclude title layout (0) if there are others
                valid_txt = txt_layouts[1:] if len(txt_layouts) > 1 else txt_layouts
                l_to_use = random.choice(valid_txt)
            else:
                l_to_use = prs.slide_layouts[0]
                
            slide = prs.slides.add_slide(l_to_use)
            
            title_injected = False
            bullets_injected = False
            
            # Clear hardcoded raw text from Gamma that isn't a placeholder
            for shape in list(slide.shapes):
                if not shape.is_placeholder and hasattr(shape, "text") and shape.has_text_frame:
                    if shape.text.strip():
                        # Destroy the hardcoded text shape to stop overlapping
                        sp = shape.element
                        sp.getparent().remove(sp)

            # Fill the placeholders dynamically
            for ph in slide.placeholders:
                if ph.placeholder_format.type == 18:
                    if has_img:
                        img_path = s.get("image_path")
                        if img_path and os.path.exists(img_path):
                            try:
                                ph.insert_picture(img_path)
                            except Exception: pass
                else:
                    if not ph.text:
                        title_str = s.get("title", "")
                        bullets_str = "\n".join([f"  {b}" for b in content])
                        
                        text_phs = [p for p in slide.placeholders if p.placeholder_format.type != 18]
                        if len(text_phs) == 1:
                            ph.text = title_str + "\n\n" + bullets_str
                        else:
                            min_idx = min(p.placeholder_format.idx for p in text_phs)
                            if ph.placeholder_format.idx == min_idx:
                                if not title_injected:
                                    ph.text = title_str
                                    title_injected = True
                            else:
                                if not bullets_injected:
                                    ph.text = bullets_str
                                    bullets_injected = True
                                else:
                                    ph.text = ""

                        # Force format the font so it isn't black
                        if ph.has_text_frame:
                            for paragraph in ph.text_frame.paragraphs:
                                paragraph.alignment = align
                                for run in paragraph.runs:
                                    run.font.color.rgb = theme["text"]


        else:
            slide = prs.slides.add_slide(blank_layout)
            set_bg(slide, theme["bg"])
            
            if style == "corporate":
                sb_x = Inches(12.88) if lang == "kurdish" else Inches(0)
                tb_x = Inches(0.5) if lang == "kurdish" else Inches(0.75)
                sidebar = slide.shapes.add_shape(1, sb_x, Inches(0), Inches(0.45), Inches(7.5))
                sidebar.fill.solid(); sidebar.fill.fore_color.rgb = theme["accent"]; sidebar.line.fill.background()
                tb = slide.shapes.add_textbox(tb_x, Inches(0.25), Inches(11.8), Inches(0.75))
                add_text(tb.text_frame, s.get("title", ""), 26, True, theme["title"], align=align)
                line_x = Inches(0.5) if lang == "kurdish" else Inches(0.75)
                line = slide.shapes.add_shape(1, line_x, Inches(1.1), Inches(11.8), Inches(0.04))
                line.fill.solid(); line.fill.fore_color.rgb = theme["accent"]; line.line.fill.background()
                text_x = Inches(0.5) if lang == "kurdish" else Inches(0.75)
                text_y, text_w, text_h = Inches(1.25), Inches(7.8) if has_img else Inches(11.8), Inches(5.8)
                bullet_size = 17

            elif style == "bold":
                hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
                hdr.fill.solid(); hdr.fill.fore_color.rgb = theme["accent"]; hdr.line.fill.background()
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.9))
                add_text(tb.text_frame, s.get("title", ""), 30, True, theme["bg"], align=PP_ALIGN.CENTER)
                text_x = Inches(0.5) if lang == "kurdish" and has_img else Inches(5.5) if has_img else Inches(0.8)
                text_y, text_w, text_h = Inches(1.3), Inches(7.0) if has_img else Inches(11.73), Inches(5.6)
                bullet_size = 19

            elif style == "minimal":
                strip_x = Inches(13.0) if lang == "kurdish" else Inches(0.3)
                tb_x = Inches(0.6) if lang == "kurdish" else Inches(0.6)
                strip = slide.shapes.add_shape(1, strip_x, Inches(0.3), Inches(0.04), Inches(6.9))
                strip.fill.solid(); strip.fill.fore_color.rgb = theme["accent"]; strip.line.fill.background()
                tb = slide.shapes.add_textbox(tb_x, Inches(0.25), Inches(12.0), Inches(0.75))
                add_text(tb.text_frame, s.get("title", ""), 26, True, theme["title"], align=align)
                text_x = Inches(0.6) if lang == "kurdish" and has_img else Inches(5.3) if has_img else Inches(0.6)
                text_y, text_w, text_h = Inches(1.2), Inches(7.4) if has_img else Inches(12.0), Inches(5.8)
                bullet_size = 17

            elif style == "tech":
                top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.1))
                top.fill.solid(); top.fill.fore_color.rgb = theme["accent"]; top.line.fill.background()
                bot = slide.shapes.add_shape(1, Inches(0), Inches(7.4), Inches(13.33), Inches(0.1))
                bot.fill.solid(); bot.fill.fore_color.rgb = theme["accent"]; bot.line.fill.background()
                for cx, cy in [(Inches(0.2), Inches(0.2)), (Inches(12.9), Inches(7.1))]:
                    sq = slide.shapes.add_shape(1, cx, cy, Inches(0.14), Inches(0.14))
                    sq.fill.solid(); sq.fill.fore_color.rgb = theme["accent"]; sq.line.fill.background()
                tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.73), Inches(0.75))
                add_text(tb.text_frame, s.get("title", ""), 27, True, theme["title"], align=PP_ALIGN.CENTER)
                div = slide.shapes.add_shape(1, Inches(0.8), Inches(1.05), Inches(11.73), Inches(0.04))
                div.fill.solid(); div.fill.fore_color.rgb = theme["accent"]; div.line.fill.background()
                text_x = Inches(0.8) if lang == "kurdish" and has_img else Inches(5.3) if has_img else Inches(0.8)
                text_y, text_w, text_h = Inches(1.2), Inches(7.4) if has_img else Inches(11.73), Inches(5.8)
                bullet_size = 18

            elif style == "elegant":
                sb_x = Inches(0) if lang == "kurdish" else Inches(12.88)
                tb_x = Inches(0.8) if lang == "kurdish" else Inches(0.5)
                sidebar = slide.shapes.add_shape(1, sb_x, Inches(0), Inches(0.45), Inches(7.5))
                sidebar.fill.solid(); sidebar.fill.fore_color.rgb = theme["accent"]; sidebar.line.fill.background()
                tb = slide.shapes.add_textbox(tb_x, Inches(0.25), Inches(12.0), Inches(0.75))
                add_text(tb.text_frame, s.get("title", ""), 27, True, theme["title"], align=align)
                line_x = Inches(0.8) if lang == "kurdish" else Inches(7.5)
                line = slide.shapes.add_shape(1, line_x, Inches(1.1), Inches(5.0), Inches(0.04))
                line.fill.solid(); line.fill.fore_color.rgb = theme["accent"]; line.line.fill.background()
                text_x = Inches(0.8) if lang == "kurdish" and has_img else Inches(5.3) if has_img else Inches(0.5)
                text_y, text_w, text_h = Inches(1.25), Inches(7.0) if has_img else Inches(12.0), Inches(5.8)
                bullet_size = 17

            else:
                # Header block
                header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
                header.fill.solid(); header.fill.fore_color.rgb = theme["accent"]; header.line.fill.background()
                # Accent bottom strip
                strip = slide.shapes.add_shape(1, Inches(0), Inches(7.35), Inches(13.33), Inches(0.15))
                strip.fill.solid(); strip.fill.fore_color.rgb = theme["accent"]; strip.line.fill.background()
                # Title inside header
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.33), Inches(1.1))
                add_text(tb.text_frame, s.get("title", ""), 30, True, theme["bg"], align=align)
                text_x = Inches(0.8) if lang == "kurdish" and has_img else Inches(5.3) if has_img else Inches(0.8)
                text_y, text_w, text_h = Inches(1.45), Inches(7.0) if has_img else Inches(11.73), Inches(5.6)
                bullet_size = 20

            tb_content = slide.shapes.add_textbox(text_x, text_y, text_w, text_h)
            tf = tb_content.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for i, bullet in enumerate(content):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.alignment = align; p.space_after = Pt(12)
                run = p.add_run()
                run.text = f"●  {bullet}"
                run.font.size = Pt(bullet_size)
                run.font.color.rgb = theme["text"]

            if has_img:
                img_x = (text_x + text_w + Inches(0.2)) if lang == "kurdish" else Inches(0.5)
                if style == "corporate":
                    img_x = Inches(0.5) if lang == "kurdish" else Inches(8.8)
                _add_image(slide, s, img_x, Inches(1.5), Inches(4.4), Inches(3.6))

            num_x = Inches(12.8) if style in ("corporate", "classic", "bold", "minimal", "tech") else Inches(0.3)
            tb = slide.shapes.add_textbox(num_x, Inches(6.95), Inches(0.8), Inches(0.4))
            add_text(tb.text_frame, str(idx + 2), 12, False, theme["accent"])

    prs.save(output_path)
